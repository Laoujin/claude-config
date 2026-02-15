from pptx import Presentation
import sys, io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

if len(sys.argv) < 2:
    print("Usage: python extract_pptx.py <path-to-pptx>")
    sys.exit(1)

prs = Presentation(sys.argv[1])
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
    print()
